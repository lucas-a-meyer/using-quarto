from pptx.util import Pt
from pptx import Presentation

prs = Presentation("/home/lucasmeyer/msft/using-quarto/docs/presentation.pptx")

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.name == "Consolas":
                    run.font.size = Pt(11)

prs.save("/home/lucasmeyer/msft/using-quarto/docs/presentation-fixed.pptx")
